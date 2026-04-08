#!/usr/bin/env python3
"""
특허제안서 PPT 자동 생성 스크립트
사용법: python3 proposal/fill_template.py proposal/input_data.json
출력:   proposal/output_YYMMDD_HHMMSS.pptx
"""

import json
import sys
import os
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'proposal_template.pptx')
IMAGES_DIR = os.path.join(os.path.dirname(__file__), 'images')

# 포트폴리오 이미지 → shape name 매핑
# 슬라이드1 기준 이미지 shape 순서 (원본 분석 결과)
IMAGE_SHAPE_MAP = {
    'img_1.jpg': 'object 17',   # 행1 좌
    'img_2.jpg': 'object 23',   # 행1 중
    'img_3.jpg': 'object 26',   # 행1 우  (※ 원본에서 object 26이 (2.96in, 9.90in)이므로 순서 확인 필요)
    'img_4.jpg': 'object 29',   # 행2 좌
    'img_5.jpg': 'object 32',   # 행2 중
    'img_6.jpg': 'object 20',   # 행2 우
    'img_7.jpg': 'object 41',   # 행3 좌
    'img_8.jpg': 'object 44',   # 행3 중
    'img_9.jpg': 'object 47',   # 행3 우
}

def fill_text(prs, data: dict):
    """모든 슬라이드의 placeholder를 data 딕셔너리로 치환"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, value in data.items():
                            placeholder = '{{' + key + '}}'
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, str(value))

        # 테이블 내 placeholder 치환
        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                for key, value in data.items():
                                    placeholder = '{{' + key + '}}'
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, str(value))


def replace_image(slide, shape_name: str, image_path: str):
    """지정 shape의 이미지를 새 이미지로 교체"""
    from pptx.oxml.ns import qn
    from lxml import etree
    import base64

    for shape in slide.shapes:
        if shape.name == shape_name and shape.shape_type == 13:  # PICTURE
            # 이미지 바이너리 교체
            pic = shape._element
            # 기존 blip 찾기
            blip = pic.find('.//' + qn('a:blip'))
            if blip is None:
                print(f"  ⚠️  {shape_name}: blip 요소를 찾을 수 없음 (수동 교체 필요)")
                return

            # rId 확인
            embed = blip.get(qn('r:embed'))
            if embed:
                # 관련 이미지 파트 교체
                img_part = slide.part.related_parts[embed]
                with open(image_path, 'rb') as f:
                    img_data = f.read()
                img_part._blob = img_data
                print(f"  ✅ {shape_name} → {os.path.basename(image_path)} 교체 완료")
            return

    print(f"  ⚠️  {shape_name} shape을 찾을 수 없음")


def main():
    if len(sys.argv) < 2:
        print("사용법: python3 proposal/fill_template.py <input_data.json>")
        sys.exit(1)

    json_path = sys.argv[1]
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    print(f"입력 파일: {json_path}")
    print(f"템플릿: {TEMPLATE_PATH}")

    # 출력 파일명
    timestamp = datetime.now().strftime('%y%m%d_%H%M%S')
    output_path = os.path.join(os.path.dirname(__file__), f'output_{timestamp}.pptx')

    # 템플릿 복사
    shutil.copy(TEMPLATE_PATH, output_path)
    prs = Presentation(output_path)

    # 1) 텍스트 치환
    print("\n[1/2] 텍스트 치환 중...")
    fill_text(prs, data)

    # 2) 이미지 교체
    print("\n[2/2] 이미지 교체 중...")
    slide1 = prs.slides[0]
    for img_filename, shape_name in IMAGE_SHAPE_MAP.items():
        img_path = os.path.join(IMAGES_DIR, img_filename)
        if os.path.exists(img_path):
            replace_image(slide1, shape_name, img_path)
        else:
            print(f"  ⚠️  {img_filename} 없음 (건너뜀)")

    prs.save(output_path)
    print(f"\n✅ 완료: {output_path}")
    print("\n⚠️  수동 확인 필요 항목:")
    print("  - 이미지 비율/위치 확인")
    print("  - 포트폴리오 링크 버튼 하이퍼링크 설정 (도형 우클릭 → 하이퍼링크 편집)")
    print("  - 사업수행 기간 표 행 수 (9행 초과/미달 시 수동 수정)")
    print("  - 역할 표 연구원 수 변동 시 수동 수정")


if __name__ == '__main__':
    main()
